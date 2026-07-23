import sys
import streamlit as st
import streamlit.components.v1 as components
from google import genai
from google.genai import types
import pandas as pd
import json
import re
import time
from pptx import Presentation
import docx
import plotly.graph_objects as go

# --- 1. アプリ設定 & UI ---
st.set_page_config(page_title="AI Analyzer", layout="wide", initial_sidebar_state="expanded")

# 画面の余白最適化 ＆ PDF出力（印刷）用のカスタムCSS
st.html("""
    <style>
        [data-testid="stSidebarHeader"] {
            padding-top: 0.5rem !important;
            padding-bottom: 0rem !important;
            min-height: auto !important;
        }
        [data-testid="stSidebarUserContent"] {
            padding-top: 0rem !important;
        }
        .block-container {
            padding-top: 3.5rem !important; 
            padding-bottom: 2rem !important;
        }
        
        /* 🖨️ PDF出力（印刷）時専用のスタイル */
        @media print {
            [data-testid="stSidebar"] { display: none !important; }
            header[data-testid="stHeader"] { display: none !important; }
            .no-print { display: none !important; }
            iframe { display: none !important; }
            .block-container { padding-top: 0rem !important; padding-bottom: 0rem !important; max-width: 100% !important; }
            * {
                -webkit-print-color-adjust: exact !important;
                color-adjust: exact !important;
                print-color-adjust: exact !important;
            }
            
            /* A4印刷時の横並び維持と、はみ出し防止 */
            [data-testid="stHorizontalBlock"] {
                display: flex !important;
                flex-direction: row !important;
                flex-wrap: nowrap !important;
                align-items: flex-start !important;
            }
            [data-testid="column"] {
                min-width: 0 !important;
            }
            
            /* Plotly（画像化されたグラフ）を枠内に強制収容 */
            .js-plotly-plot, .plotly, .plotly svg {
                max-width: 100% !important;
                height: auto !important;
            }
        }
    </style>
""")

# タイトルエリア
st.html("""
    <div style="background-color:#f8f9fa; padding:20px; border-radius:10px; border-left:8px solid #007bff; margin-bottom:25px;">
        <h2 style="margin:0; color:#1e293b;">📊 AI Analyzer Pro</h2>
        <p style="margin:5px 0 0 0; color:#64748b; font-size:15px;">自社のブランド戦略が生成AIにどこまで正しく伝わっているかを分析し、次の一手を見つけるダッシュボード</p>
    </div>
""")

# セッション状態の初期化
if "bas_result" not in st.session_state:
    st.session_state.bas_result = None

# サイドバー：設定エリア
with st.sidebar:
    st.markdown("### ⚙️ 解析設定")
    st.divider()
    
    api_key = st.text_input("🔑 Gemini API Key", type="password", help="Google AI Studioで発行したAPIキーを入力してください。")
    brand_name = st.text_input("🏷️ ブランド名", value="ディアナチュラ")
    brand_url = st.text_input("🌐 公式サイトURL", value="https://www.dear-natura.com/")
    
    st.divider()
    st.markdown("### 📂 データのアップロード")
    uploaded_pptx = st.file_uploader("❶ 生成AI分析資料（PPTX）", type=["pptx"])
    uploaded_csv = st.file_uploader("❷ 生成AIでの言及数データ(CSV)", type=["csv", "txt"])
    uploaded_docx = st.file_uploader("❸ 生成AIのブランド評価(DOCX)", type=["docx"])
    
    st.divider()
    if st.button("⏹️ 処理を中断 / 画面をリセット", type="secondary", use_container_width=True):
        st.session_state.bas_result = None
        st.rerun()
        
    st.divider()
    debug_mode = st.checkbox("🔧 開発者用ログを表示する", value=False)


# --- ユーティリティ関数 ---
def clean_and_parse_json(text: str) -> dict:
    json_match = re.search(r'\{.*\}', text, re.DOTALL)
    if not json_match:
        raise ValueError("JSONが見つかりません")
    text = json_match.group(0)
    text = re.sub(r'[\x00-\x1F\x7F]', '', text)
    return json.loads(text)

def extract_text_from_pptx(file) -> str:
    try:
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)[:3000]
    except Exception as e:
        return f"PPTXエラー: {e}"

def extract_text_from_docx(file) -> str:
    try:
        doc = docx.Document(file)
        text = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        return "\n".join(text)[:3000]
    except Exception as e:
        return f"DOCXエラー: {e}"

def load_csv_data(file):
    encodings = ["utf-8", "shift_jis", "cp932", "utf-8-sig", "iso-8859-1"]
    delimiters = [",", "\t", ";"]
    for encoding in encodings:
        for delimiter in delimiters:
            try:
                file.seek(0)
                df = pd.read_csv(file, encoding=encoding, sep=delimiter, engine="python", on_bad_lines="skip")
                if not df.empty and len(df.columns) >= 2:
                    return df
            except Exception:
                continue
    return None


# --- 2. 分析ロジック ---
if st.button("🚀 戦略ギャップ分析を実行", type="primary", use_container_width=True):
    if not api_key or not brand_url or not uploaded_pptx or not uploaded_csv or not uploaded_docx:
        st.error("左側のサイドバーで、APIキー、URL、および3つのファイルをすべてセットしてください。")
        st.stop()

    with st.spinner("AIがデータを解析し、経営・マーケティング向けのレポートを作成中..."):
        try:
            client = genai.Client(api_key=api_key)
            model_name = "gemini-2.5-flash"
            
            pptx_text = extract_text_from_pptx(uploaded_pptx)
            df_raw = load_csv_data(uploaded_csv)
            docx_text = extract_text_from_docx(uploaded_docx)
            
            if df_raw is None:
                st.error("CSVファイルの読み込みに失敗しました。")
                st.stop()
                
            csv_context = df_raw.head(35).to_csv(index=False)
            
            # Phase 1: オウンドメディアからのキーワード抽出
            prompt_dict = f"""
            ブランド "{brand_name}" (公式URL: {brand_url}) の戦略資料からオウンドメディアが目指すブランド像を示すキーワードを5つずつ抽出してください。
            
            [戦略資料テキスト]
            {pptx_text}
            """
            
            dict_schema = {
                "type": "object",
                "properties": {
                    "core": {"type": "array", "items": {"type": "string"}},
                    "functional": {"type": "array", "items": {"type": "string"}},
                    "professional": {"type": "array", "items": {"type": "string"}}
                },
                "required": ["core", "functional", "professional"]
            }

            dictionary_data = None
            for attempt in range(3):
                try:
                    res_dict = client.models.generate_content(
                        model=model_name,
                        contents=[prompt_dict],
                        config=types.GenerateContentConfig(temperature=0.1, response_mime_type="application/json", response_schema=dict_schema)
                    )
                    dictionary_data = json.loads(res_dict.text)
                    break
                except Exception as e:
                    if "503" in str(e) and attempt < 2:
                        time.sleep(3)
                        continue
                    raise e

            if not dictionary_data:
                st.error("AIサーバーが混雑しています。少し時間を置いて再度お試しください。")
                st.stop()

            # Phase 2: 戦略的ギャップ分析とスコアリング
            response_schema = {
                "type": "object",
                "properties": {
                    "diagnosis_story": {
                        "type": "object",
                        "properties": {
                            "match": {"type": "string"},
                            "positive_gap": {"type": "string"},
                            "negative_gap": {"type": "string"}
                        },
                        "required": ["match", "positive_gap", "negative_gap"]
                    },
                    "topline": {"type": "string"},
                    "competitive_analysis": {
                        "type": "object",
                        "properties": {
                            "benchmark_competitors": {"type": "string"},
                            "mention_volume_comparison": {"type": "string"},
                            "mention_order_comparison": {"type": "string"},
                            "mention_content_comparison": {"type": "string"},
                            "strategic_advice": {"type": "string"}
                        },
                        "required": ["benchmark_competitors", "mention_volume_comparison", "mention_order_comparison", "mention_content_comparison", "strategic_advice"]
                    },
                    "improvement_actions": {"type": "array", "items": {"type": "string"}},
                    "detailed_discrepancies": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "issue": {"type": "string"},
                                "impact": {"type": "string"},
                                "solution": {"type": "string"}
                            },
                            "required": ["issue", "impact", "solution"]
                        }
                    },
                    "radar_quantity": {
                        "type": "object",
                        "properties": {
                            "brand_philosophy": {"type": "integer"},
                            "functional_value": {"type": "integer"},
                            "emotional_engagement": {"type": "integer"},
                            "safety_reputation": {"type": "integer"},
                            "usage_scene_moment": {"type": "integer"}
                        },
                        "required": ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
                    },
                    "radar_quantity_summary": {"type": "string"},
                    "radar_quality": {
                        "type": "object",
                        "properties": {
                            "brand_philosophy": {"type": "integer"},
                            "functional_value": {"type": "integer"},
                            "emotional_engagement": {"type": "integer"},
                            "safety_reputation": {"type": "integer"},
                            "usage_scene_moment": {"type": "integer"}
                        },
                        "required": ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
                    },
                    "radar_quality_summary": {"type": "string"},
                    "radar_reasons": {
                        "type": "object",
                        "properties": {
                            "brand_philosophy": {"type": "string"},
                            "functional_value": {"type": "string"},
                            "emotional_engagement": {"type": "string"},
                            "safety_reputation": {"type": "string"},
                            "usage_scene_moment": {"type": "string"}
                        },
                        "required": ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
                    }
                },
                "required": ["diagnosis_story", "topline", "competitive_analysis", "improvement_actions", "detailed_discrepancies", "radar_quantity", "radar_quantity_summary", "radar_quality", "radar_quality_summary", "radar_reasons"]
            }

            prompt_analysis = f"""
            Analyze Generative AI output data for "{brand_name}" (Official URL: {brand_url}) against the owned media keywords.
            
            [OWNED MEDIA KEYWORDS]
            {json.dumps(dictionary_data, ensure_ascii=False)}
            
            [GENERATIVE AI RANKING DATA (CSV)]
            {csv_context}
            
            [GENERATIVE AI BRAND EVALUATION (DOCX)]
            {docx_text}
            
            TASK:
            1. "diagnosis_story": Write 3 fluent narrative paragraphs in Japanese (EACH strictly around 200-250 characters) aimed at business executives.
               - "match": What aspects of the company's message are correctly understood by the AI?
               - "positive_gap": What unexpected strengths or positive perceptions did the AI find?
               - "negative_gap": What misconceptions or weak points exist in the AI's understanding?
            2. "topline": Write a single-sentence summary strategy for executives.
            3. "competitive_analysis": Analyze how the AI perceives the brand relative to its competitors based ONLY on the provided data. DO NOT USE NUMERICAL SCORES AND DO NOT LIST SIMPLE WIN/LOSS.
               - "benchmark_competitors": ベンチマークとすべき企業 (Qualitatively suggest 1-2 competitor brands that this brand should benchmark against in AI recommendations based on mention frequency, order, and context. Explain WHY qualitatively. Around 80-120 characters).
               - "mention_volume_comparison": 言及の多さの比較 (Compare the frequency/volume of mentions of this brand versus competitors in the data. Around 80-100 characters).
               - "mention_order_comparison": 言及順番の比較 (Compare the ranking or order in which this brand is mentioned versus competitors. Are competitors mentioned first? Around 80-100 characters).
               - "mention_content_comparison": 言及内容の比較 (Compare the qualitative content/context of mentions. How does the AI describe this brand compared to competitors? Around 80-100 characters).
               - "strategic_advice": Provide one highly specific strategic action to improve the brand's position against the benchmarked competitors.
            4. "improvement_actions": Provide EXACTLY 5 clear, actionable marketing steps.
            5. "detailed_discrepancies": Identify up to 10 HIGHLY SPECIFIC perception issues or missing elements in the Generative AI's understanding. 
               CRITICAL INSTRUCTION: Do NOT explicitly assert or guess the company's intended message. Focus entirely on what the AI currently outputs. Every item MUST explicitly quote specific data points, quotes, or ranks from [GENERATIVE AI RANKING DATA] or [GENERATIVE AI BRAND EVALUATION].
               - "issue": Detail the specific AI perception issue based ONLY on the provided AI data.
               - "impact": Explain the specific business impact tailored to THIS brand's actual product and market.
               - "solution": Provide a concrete, highly specific PR/Marketing action to fix this AI perception gap.
            6. "radar_quantity", "radar_quality", summaries & "radar_reasons": Score the Generative AI's perception in PERCENTAGE (0-100) for the following 5 criteria from TWO perspectives:
               - "radar_quantity" (量的乖離/一致確率): Estimate the % probability (0-100) that the AI's answer MATCHES the owned media.
               - "radar_quantity_summary": Write a brief overview (approx. 100-150 characters in Japanese) summarizing the overall shape of the quantitative radar chart.
               - "radar_quality" (質的乖離/類似度): Estimate the % similarity (0-100) of the AI's answers compared to the owned media.
               - "radar_quality_summary": Write a brief overview (approx. 100-150 characters in Japanese) summarizing the overall shape of the qualitative radar chart.
               CRITICAL for "radar_reasons": Provide a DETAILED business reason explaining BOTH the quantity and quality scores based on the data.
               Criteria:
               - "brand_philosophy": ブランド理念
               - "functional_value": 機能価値
               - "emotional_engagement": 情緒的エンゲージメント
               - "safety_reputation": 安全性と評判
               - "usage_scene_moment": 利用シーン・モーメント一致度 (Match in usage context/timing). CRITICAL INSTRUCTION: Do NOT score 0% simply because the owned media keywords omit specific usage scenes. If the AI's suggested usage broadly aligns with the common sense/expected usage of this product type (e.g., 'daily use' vs 'only when tired'), consider it a match and score it favorably.
            Return JSON in Japanese.
            """

            final_data = None
            for attempt in range(3):
                try:
                    res_analysis = client.models.generate_content(
                        model=model_name,
                        contents=[prompt_analysis],
                        config=types.GenerateContentConfig(temperature=0.1, response_mime_type="application/json", response_schema=response_schema)
                    )
                    final_data = json.loads(res_analysis.text)
                    break
                except Exception as e:
                    if "503" in str(e) and attempt < 2:
                        time.sleep(3)
                        continue
                    raise e

            if final_data:
                final_data["dictionary"] = dictionary_data
                st.session_state.bas_result = final_data
            else:
                st.error("AIサーバーが混雑しています。少し時間を置いて再度お試しください。")
                st.stop()

        except Exception as e:
            st.error(f"分析エラーが発生しました: {e}")
            st.stop()


# --- 3. 結果表示 UI ---
if st.session_state.bas_result:
    res = st.session_state.bas_result

    # 🖨️ PDF出力ボタンの配置
    components.html("""
        <div style="display: flex; justify-content: flex-end; padding-right: 10px;">
            <button onclick="window.parent.print()" style="background-color: #475569; color: white; border: none; padding: 10px 20px; border-radius: 5px; cursor: pointer; font-weight: bold; font-size: 14px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); font-family: sans-serif;">
                📄 レポートをPDFで出力する
            </button>
        </div>
    """, height=60)

    if debug_mode and res:
        with st.expander("🔧 開発者用デバッグログ"):
            st.json(res)

    # ==========================================
    # ① 現状診断：自社発信とAIの認識ギャップ
    # ==========================================
    st.markdown("### 🎯 ① 現状の診断：自社発信とAIの認識ギャップ")
    st.caption("自社で発信しているメッセージがAIにどう伝わっているか、「狙い通りな点」と「意外なズレ」をわかりやすく解説します。")
    
    story = res.get("diagnosis_story", {})
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.html("<div style='background-color:#e6f4ea; padding:12px; border-radius:5px; border-left:5px solid #28a745; font-weight:bold; color:#137333; margin-bottom:10px;'>🟢 狙い通りに伝わっている強み</div>")
        st.write(story.get("match", "分析データなし"))
        
    with col2:
        st.html("<div style='background-color:#e8f0fe; padding:12px; border-radius:5px; border-left:5px solid #007bff; font-weight:bold; color:#1a73e8; margin-bottom:10px;'>🔵 AIが見つけた「意外な強み」</div>")
        st.write(story.get("positive_gap", "分析データなし"))
        
    with col3:
        st.html("<div style='background-color:#fce8e6; padding:12px; border-radius:5px; border-left:5px solid #dc3545; font-weight:bold; color:#c5221f; margin-bottom:10px;'>🔴 AIの誤解や情報不足（課題）</div>")
        st.write(story.get("negative_gap", "分析データなし"))
            
    st.divider()

    # ==========================================
    # ② 生成AIからのブランド評価（量的乖離・質的乖離）
    # ==========================================
    st.markdown("### 📊 ② 生成AIからのブランド評価（2軸による乖離分析）")
    st.caption("AIの認識ズレを「どれくらいの頻度で一致するか（量的乖離）」と「自社発信とどれくらい内容が似ているか（質的乖離）」の2軸（単位：％）で可視化しています。")
    
    q_qty = res.get("radar_quantity", {})
    q_qual = res.get("radar_quality", {})
    qty_summary = res.get("radar_quantity_summary", "サマリーデータがありません。")
    qual_summary = res.get("radar_quality_summary", "サマリーデータがありません。")
    reasons = res.get("radar_reasons", {})
    
    categories = ['ブランド理念', '機能価値', '情緒的<br>エンゲージメント', '安全性と評判', '利用シーン・<br>モーメント']
    categories_closed = categories + [categories[0]]
    
    keys = ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
    qty_scores = [q_qty.get(k, 0) for k in keys]
    qual_scores = [q_qual.get(k, 0) for k in keys]
    
    qty_closed = qty_scores + [qty_scores[0]]
    qual_closed = qual_scores + [qual_scores[0]]

    # --- 1段目：1. 量的乖離 ---
    st.markdown("#### 🔵 1. 量的乖離（自社発信と一致する確率：％）")
    col_chart1, col_summary1 = st.columns([1, 1.2])
    
    with col_chart1:
        fig_qty = go.Figure()
        fig_qty.add_trace(go.Scatterpolar(
            r=qty_closed,
            theta=categories_closed,
            fill='toself',
            name='量的乖離',
            line_color='#1a73e8',
            fillcolor='rgba(26, 115, 232, 0.2)'
        ))
        fig_qty.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 100], ticksuffix="%")),
            showlegend=False,
            margin=dict(l=40, r=40, t=30, b=30),
            height=320 
        )
        st.plotly_chart(fig_qty, use_container_width=True, config={'staticPlot': True})
        
    with col_summary1:
        st.html(f"""
        <div style="margin-top: 20px; padding: 25px; background-color: #f8f9fa; border-left: 6px solid #1a73e8; border-radius: 6px; box-shadow: 0 2px 4px rgba(0,0,0,0.05);">
            <div style="font-weight: bold; color: #1a73e8; margin-bottom: 12px; font-size: 16px;">
                📈 量的乖離の全体傾向
            </div>
            <div style="font-size: 15px; color: #333; line-height: 1.8;">
                {qty_summary}
            </div>
        </div>
        """)

    st.markdown("<br>", unsafe_allow_html=True)

    # --- 2段目：2. 質的乖離 ---
    st.markdown("#### 🟢 2. 質的乖離（自社発信と内容の類似度：％）")
    col_chart2, col_summary2 = st.columns([1, 1.2])
    
    with col_chart2:
        fig_qual = go.Figure()
        fig_qual.add_trace(go.Scatterpolar(
            r=qual_closed,
            theta=categories_closed,
            fill='toself',
            name='質的乖離',
            line_color='#28a745',
            fillcolor='rgba(40, 167, 69, 0.2)'
        ))
        fig_qual.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 100], ticksuffix="%")),
            showlegend=False,
            margin=dict(l=40, r=40, t=30, b=30),
            height=320
        )
        st.plotly_chart(fig_qual, use_container_width=True, config={'staticPlot': True})
        
    with col_summary2:
        st.html(f"""
        <div style="margin-top: 20px; padding: 25px; background-color: #f8f9fa; border-left: 6px solid #28a745; border-radius: 6px; box-shadow: 0 2px 4px rgba(0,0,0,0.05);">
            <div style="font-weight: bold; color: #28a745; margin-bottom: 12px; font-size: 16px;">
                📈 質的乖離の全体傾向
            </div>
            <div style="font-size: 15px; color: #333; line-height: 1.8;">
                {qual_summary}
            </div>
        </div>
        """)

    st.divider()

    # 評価理由をカード型デザインで表示
    st.markdown("#### 📝 各項目の評価詳細（なぜこの数値になったのか）")
    
    for title, key in zip(['ブランド理念の浸透度', '機能価値の伝達度', '情緒的エンゲージメント', 'ブランドの安全性と評判', '利用シーン・モーメント一致度'], keys):
        qty_val = q_qty.get(key, 0)
        qual_val = q_qual.get(key, 0)
        reason = reasons.get(key, 'データなし')
        
        st.html(f"""
        <div style="border-left: 5px solid #007bff; background-color: #f8f9fa; padding: 15px; margin-bottom: 12px; border-radius: 4px; box-shadow: 0 1px 2px rgba(0,0,0,0.05);">
            <div style="font-weight: bold; font-size: 16px; margin-bottom: 8px; color: #333;">
                🔹 {title}
            </div>
            <div style="display: flex; gap: 15px; margin-bottom: 10px;">
                <div style="background-color: #e8f0fe; padding: 4px 10px; border-radius: 4px; font-size: 13px; color: #1a73e8; font-weight: bold;">
                    量的乖離（一致確率）: {qty_val}%
                </div>
                <div style="background-color: #e6f4ea; padding: 4px 10px; border-radius: 4px; font-size: 13px; color: #137333; font-weight: bold;">
                    質的乖離（内容の類似度）: {qual_val}%
                </div>
            </div>
            <div style="font-size: 14px; color: #555; line-height: 1.6;">
                {reason}
            </div>
        </div>
        """)

    st.divider()
    
    # ==========================================
    # ③ 対競合優先度の分析
    # ==========================================
    st.markdown("### ⚔️ ③ 対競合優先度の分析")
    st.caption("AIの回答データ内において、自社が競合他社と比較してどのように言及され、どのようなポジションにいるかを定性的に分析します。")
    
    comp_data = res.get("competitive_analysis", {})
    bench_comp = comp_data.get("benchmark_competitors", "データなし")
    vol_comp = comp_data.get("mention_volume_comparison", "データなし")
    order_comp = comp_data.get("mention_order_comparison", "データなし")
    content_comp = comp_data.get("mention_content_comparison", "データなし")
    
    st.html(f"""
    <div style="background-color: #fcfaff; border-left: 5px solid #8b5cf6; padding: 18px; border-radius: 4px; box-shadow: 0 1px 3px rgba(0,0,0,0.05); margin-bottom: 25px;">
        <div style="font-weight: bold; color: #7c3aed; margin-bottom: 8px; font-size: 16px;">🎯 ベンチマークとすべき競合企業とその理由</div>
        <div style="font-size: 15px; color: #333; line-height: 1.6;">{bench_comp}</div>
    </div>
    
    <div style="margin-bottom: 20px;">
        <div style="border-left: 4px solid #475569; background-color: #f8f9fa; padding: 15px; margin-bottom: 10px; border-radius: 0 4px 4px 0;">
            <div style="font-weight: bold; color: #334155; font-size: 14px; margin-bottom: 4px;">📊 言及の多さの比較</div>
            <div style="font-size: 14px; color: #333; line-height: 1.5;">{vol_comp}</div>
        </div>
        <div style="border-left: 4px solid #475569; background-color: #f8f9fa; padding: 15px; margin-bottom: 10px; border-radius: 0 4px 4px 0;">
            <div style="font-weight: bold; color: #334155; font-size: 14px; margin-bottom: 4px;">🔢 言及順番の比較（第一想起されているか）</div>
            <div style="font-size: 14px; color: #333; line-height: 1.5;">{order_comp}</div>
        </div>
        <div style="border-left: 4px solid #475569; background-color: #f8f9fa; padding: 15px; border-radius: 0 4px 4px 0;">
            <div style="font-weight: bold; color: #334155; font-size: 14px; margin-bottom: 4px;">💬 言及内容の比較（どのように語られているか）</div>
            <div style="font-size: 14px; color: #333; line-height: 1.5;">{content_comp}</div>
        </div>
    </div>
    
    <div style="background-color: #e8f0fe; border: 1px solid #c2d7fa; padding: 15px; border-radius: 8px;">
        <div style="font-weight: bold; color: #1a73e8; margin-bottom: 8px;">💡 対競合の戦略的アドバイス</div>
        <div style="font-size: 15px; color: #333; line-height: 1.6;">{comp_data.get('strategic_advice', 'データなし')}</div>
    </div>
    """)
    
    st.divider()

    # ==========================================
    # ④ 今後、マーケティングで打つべき具体策
    # ==========================================
    st.markdown("### 🚀 ④ 今後、マーケティングで打つべき具体策")
    st.caption("AIの誤解を解き、意外な強みを自社のPRに活かすための5つのアクションプランです。")
    
    st.info(f"**【戦略方針】** {res.get('topline')}")
    
    st.markdown("#### 🛠️ 今後打つべき5つのアクション")
    for i, action in enumerate(res.get("improvement_actions", []), 1):
        st.html(f"""
        <div style="background-color: #ffffff; border: 1px solid #e2e8f0; border-radius: 8px; padding: 15px; margin-bottom: 12px; display: flex; align-items: flex-start; box-shadow: 0 1px 3px rgba(0,0,0,0.05);">
            <div style="background-color: #007bff; color: white; border-radius: 50%; min-width: 28px; height: 28px; display: flex; justify-content: center; align-items: center; font-weight: bold; margin-right: 15px; flex-shrink: 0;">
                {i}
            </div>
            <div style="font-size: 15px; color: #334155; line-height: 1.5; padding-top: 2px;">
                {action}
            </div>
        </div>
        """)
        
    st.divider()

    # ==========================================
    # ⑤ 重要な乖離の詳細と解決策
    # ==========================================
    st.markdown("### 🔍 ⑤ 重要な乖離の詳細と解決策")
    st.caption("細かなデータ分析に基づき、ビジネスへの影響が大きい乖離ポイントを深掘りし、具体的な解決策の糸口を提示します。")
    
    discrepancies = res.get("detailed_discrepancies", [])
    if discrepancies:
        for item in discrepancies:
            issue = item.get("issue", "")
            impact = item.get("impact", "")
            solution = item.get("solution", "")
            
            st.html(f"""
            <div style="border: 1px solid #e2e8f0; border-radius: 8px; padding: 18px; margin-bottom: 15px; background-color: #ffffff; box-shadow: 0 1px 3px rgba(0,0,0,0.05);">
                <div style="font-weight: bold; color: #b91c1c; font-size: 16px; margin-bottom: 10px; display: flex; align-items: center;">
                    <span style="margin-right: 8px; font-size: 18px;">⚠️</span> 乖離: {issue}
                </div>
                <div style="font-size: 14px; color: #334155; margin-bottom: 12px; padding-left: 26px;">
                    <span style="font-weight: bold; color: #475569; display: block; margin-bottom: 4px;">💥 ビジネスへの影響:</span>
                    {impact}
                </div>
                <div style="font-size: 14px; color: #334155; padding-top: 12px; padding-left: 26px; border-top: 1px dashed #cbd5e1;">
                    <span style="font-weight: bold; color: #0284c7; display: block; margin-bottom: 4px;">💡 解決への具体的なアドバイス:</span>
                    {solution}
                </div>
            </div>
            """)
    else:
        st.write("重要な乖離は見つかりませんでした。")
        
    st.divider()

    # 参考情報
    with st.expander("📄 参考：AIが分析のベースにした自社の公式キーワード"):
        d = res.get("dictionary", {})
        c1, c2, c3 = st.columns(3)
        with c1:
            st.markdown("**💡 コアとなる価値**")
            st.write(", ".join(d.get("core", [])) if d.get("core") else "抽出なし")
        with c2:
            st.markdown("**✨ 機能や効能**")
            st.write(", ".join(d.get("functional", [])) if d.get("functional") else "抽出なし")
        with c3:
            st.markdown("**🛡️ 信頼感や専門性**")
            st.write(", ".join(d.get("professional", [])) if d.get("professional") else "抽出なし")

else:
    st.html("""
        <div style="text-align:center; padding:100px 20px; color:#94a3b8;">
            <p style="font-size:40px; margin:0;">📥</p>
            <h4 style="margin:10px 0 0 0; color:#64748b;">データがセットされていません</h4>
            <p style="font-size:14px; margin:5px 0 0 0;">左側のサイドバーにAPIキーを入力し、3つのファイルをセットして「分析を開始する」を押してください。</p>
        </div>
    """)
